"""
AI Features Router — Recommendations, Insights, Summarizer
"""
from fastapi import APIRouter, Depends, HTTPException
from pydantic import BaseModel
from database import get_db
from auth import get_current_user, require_role
from utils.bhashini import translate
from routers.chatbot import call_gemini_with_fallback, extract_title
import logging, json, re
from datetime import datetime, timedelta
from bson import ObjectId

router = APIRouter(prefix="/api/ai", tags=["AI Features"])
logger = logging.getLogger(__name__)


@router.get("/recommendations")
async def get_recommendations(user: dict = Depends(get_current_user)):
    if user["role"] != "employee":
        raise HTTPException(403, "Only employees can get recommendations")
    db = get_db()
    uid = str(user["_id"])
    cached = await db.ai_recommendations.find_one({"user_id": uid})
    if cached and cached.get("expires_at", datetime.min) > datetime.utcnow():
        return {"recommendations": cached.get("recommendations", []), "cached": True}
    try:
        dept = user.get("department", "Unknown")
        doj = user.get("date_of_joining")
        exp = 0
        if doj:
            try:
                exp = (datetime.utcnow() - datetime.fromisoformat(str(doj).replace("Z",""))).days // 30
            except: pass
        progress = await db.progress.find({"user_id": uid}).to_list(100)
        assigned_ids = {p["course_id"] for p in progress}
        completed = []
        for p in progress:
            if p["status"] == "completed":
                c = await db.courses.find_one({"_id": ObjectId(p["course_id"])})
                if c: completed.append(f"{extract_title(c.get('title'),'en')} ({p.get('score',0)}%)")
        all_courses = await db.courses.find({}).to_list(100)
        unassigned = []
        for c in all_courses:
            if str(c["_id"]) not in assigned_ids:
                t = extract_title(c.get("title"), "en")
                d = c.get("description","")
                if isinstance(d, dict): d = d.get("en","")
                unassigned.append(f"- {t}: {str(d)[:100]} ({c.get('training_mode','online')})")
        if not unassigned:
            return {"recommendations": [], "cached": False}
        prompt = f"""Employee: {dept} dept, {exp}mo exp. Completed: {', '.join(completed[:5]) or 'None'}.
Available courses:
{chr(10).join(unassigned[:10])}
Recommend top 3 courses. JSON only: [{{"title":"...","reason":"..."}}]"""
        payload = {"contents": [{"role":"user","parts":[{"text":prompt}]}],
                   "generationConfig": {"temperature":0.3,"maxOutputTokens":300}}
        reply = await call_gemini_with_fallback(payload)
        m = re.search(r'\[.*\]', reply, re.DOTALL)
        if not m: return {"recommendations": [], "cached": False}
        recs_raw = json.loads(m.group())
        recs = []
        for r in recs_raw[:3]:
            for c in all_courses:
                tf = c.get("title","")
                titles = list(tf.values()) if isinstance(tf,dict) else [str(tf)]
                if any(t.lower().strip()==r.get("title","").lower().strip() for t in titles if t):
                    recs.append({"course_id":str(c["_id"]),"title":extract_title(tf,"en"),
                                 "reason":r.get("reason",""),"training_mode":c.get("training_mode","online")})
                    break
        now = datetime.utcnow()
        await db.ai_recommendations.update_one({"user_id":uid},
            {"$set":{"user_id":uid,"recommendations":recs,"cached_at":now,
                     "expires_at":now+timedelta(hours=24)}}, upsert=True)
        return {"recommendations": recs, "cached": False}
    except Exception as e:
        logger.error(f"Recommendations error: {e}")
        if cached: return {"recommendations": cached.get("recommendations",[]), "cached": True}
        return {"recommendations": [], "cached": False}


@router.delete("/recommendations/cache")
async def clear_rec_cache(user: dict = Depends(get_current_user)):
    db = get_db()
    await db.ai_recommendations.delete_one({"user_id": str(user["_id"])})
    return {"message": "Cache cleared"}


async def _gen_insights(db, emp: dict, lang: str = "en"):
    uid = str(emp["_id"])
    cached = await db.ai_insights.find_one({"user_id": uid})
    if cached and cached.get("expires_at", datetime.min) > datetime.utcnow():
        ins = cached.get("insights","")
        if lang != "en" and ins:
            try: ins = await translate(ins, "en", lang)
            except: pass
        return {"insights": ins, "generated_at": cached.get("generated_at"), "cached": True}
    try:
        progress = await db.progress.find({"user_id": uid}).to_list(100)
        comp = sum(1 for p in progress if p["status"]=="completed")
        inp = sum(1 for p in progress if p["status"]=="in_progress")
        now = datetime.utcnow()
        overdue = 0
        for p in progress:
            if p["status"]!="completed" and p.get("deadline_date"):
                try:
                    if datetime.fromisoformat(str(p["deadline_date"]).replace("Z","")) < now: overdue+=1
                except: pass
        scores = []
        pairs = []
        for p in progress:
            if p["status"]=="completed" and p.get("score",0)>0:
                c = await db.courses.find_one({"_id":ObjectId(p["course_id"])})
                if c:
                    pairs.append(f"{extract_title(c.get('title'),'en')}: {p['score']}%")
                    scores.append(p["score"])
        avg = round(sum(scores)/len(scores),1) if scores else 0
        prompt = f"""Employee training data:
- Department: {emp.get('department','Unknown')}
- Completed: {comp}, In progress: {inp}, Overdue: {overdue}
- Scores: {', '.join(pairs[:5])}, Average: {avg}
Give 3 bullet points: 1 strength, 1 area to improve, 1 recommendation. Under 20 words each. Plain text bullets only."""
        payload = {"contents":[{"role":"user","parts":[{"text":prompt}]}],
                   "generationConfig":{"temperature":0.3,"maxOutputTokens":200}}
        reply = await call_gemini_with_fallback(payload)
        gen_at = datetime.utcnow()
        await db.ai_insights.update_one({"user_id":uid},
            {"$set":{"user_id":uid,"insights":reply.strip(),"generated_at":gen_at,
                     "expires_at":gen_at+timedelta(hours=12)}}, upsert=True)
        final = reply.strip()
        if lang != "en":
            try: final = await translate(final, "en", lang)
            except: pass
        return {"insights": final, "generated_at": gen_at, "cached": False}
    except Exception as e:
        logger.error(f"Insights error: {e}")
        if cached: return {"insights":cached.get("insights",""),"generated_at":cached.get("generated_at"),"cached":True}
        return {"insights":"","generated_at":None,"cached":False}


@router.get("/insights/me")
async def get_my_insights(user: dict = Depends(get_current_user)):
    if user["role"] != "employee":
        raise HTTPException(403, "Only employees")
    return await _gen_insights(get_db(), user, user.get("primary_language","en"))


@router.get("/insights/{employee_id}")
async def get_emp_insights(employee_id: str, user: dict = Depends(require_role("admin"))):
    db = get_db()
    emp = await db.users.find_one({"_id": ObjectId(employee_id)})
    if not emp: raise HTTPException(404, "Employee not found")
    return await _gen_insights(db, emp, "en")


@router.delete("/insights/cache")
async def clear_ins_cache(user: dict = Depends(get_current_user)):
    await get_db().ai_insights.delete_one({"user_id": str(user["_id"])})
    return {"message": "Cache cleared"}


from typing import Optional

class SummarizeRequest(BaseModel):
    file_id: str
    lang: Optional[str] = None

@router.post("/summarize-material")
async def summarize_material(req: SummarizeRequest, user: dict = Depends(get_current_user)):
    db = get_db()
    lang = req.lang or user.get("primary_language", "en")
    cache_key = f"{req.file_id}_{lang}"
    cached = await db.ai_summaries.find_one({"cache_key": cache_key})
    if cached and cached.get("expires_at", datetime.min) > datetime.utcnow():
        return {"summary":cached.get("summary",[]),"filename":cached.get("filename",""),"cached":True}
    try:
        from motor.motor_asyncio import AsyncIOMotorGridFSBucket
        bucket = AsyncIOMotorGridFSBucket(db, bucket_name="course_materials")
        fid = ObjectId(req.file_id)
        file_doc = await db["course_materials.files"].find_one({"_id": fid})
        if not file_doc: raise HTTPException(404, "File not found")
        fname = file_doc.get("filename","unknown")
        ctype = file_doc.get("content_type","")
        cid = file_doc.get("metadata",{}).get("course_id")
        is_pdf = fname.lower().endswith(".pdf") or "pdf" in ctype.lower()
        is_pptx = fname.lower().endswith(".pptx") or "presentation" in ctype.lower()
        if not is_pdf and not is_pptx:
            raise HTTPException(400, "Only PDF and PPTX can be summarized")
        grid_out = await bucket.open_download_stream(fid)
        file_bytes = await grid_out.read()
        text = ""
        if is_pdf:
            import fitz
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            for page in doc: text += page.get_text()[:3000]
            doc.close()
        elif is_pptx:
            from pptx import Presentation
            import io
            prs = Presentation(io.BytesIO(file_bytes))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape,"text"): text += shape.text + "\n"
        text = text[:3000].strip()
        if not text:
            return {"summary":["No text content could be extracted."],"filename":fname,"cached":False}
        prompt = f"""Summarize this training material in 5 bullet points. Under 25 words each.
Focus on key learning points.

Content:
{text}"""
        payload = {"contents":[{"role":"user","parts":[{"text":prompt}]}],
                   "generationConfig":{"temperature":0.3,"maxOutputTokens":300}}
        reply = await call_gemini_with_fallback(payload)
        summary = [l.strip().lstrip("-•*0123456789.) ").strip() for l in reply.strip().split("\n") if l.strip()][:5]
        if lang != "en":
            try: summary = [await translate(b,"en",lang) for b in summary]
            except: pass
        now = datetime.utcnow()
        await db.ai_summaries.update_one({"cache_key":cache_key},
            {"$set":{"cache_key":cache_key,"file_id":req.file_id,"lang":lang,
                     "summary":summary,"filename":fname,"cached_at":now,
                     "expires_at":now+timedelta(days=7)}}, upsert=True)
        return {"summary":summary,"filename":fname,"cached":False}
    except HTTPException: raise
    except Exception as e:
        logger.error(f"Summarize error: {e}")
        return {"summary":["Unable to generate summary."],"filename":"","cached":False}
